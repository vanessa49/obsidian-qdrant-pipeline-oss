"""
格式转换模块 - 将各种格式的文档转换为 Markdown
支持：PPTX（含图片型）、PDF（含扫描件）、DOCX、HTML、TXT、MD、TEX
"""

import base64
import io
import subprocess
from pathlib import Path
from typing import Optional

import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from llm_client import call_vision

try:
    from markitdown import MarkItDown
except ImportError:
    print("警告: markitdown 未安装，部分功能可能不可用")
    MarkItDown = None


class ConversionError(Exception):
    """转换错误基类"""
    pass


class PPTXConversionError(ConversionError):
    """PPTX 转换错误"""
    pass


class PDFConversionError(ConversionError):
    """PDF 转换错误"""
    pass


def _image_to_base64(image: Image.Image) -> str:
    """将 PIL Image 转换为 base64 字符串"""
    buffered = io.BytesIO()
    # 转换为 RGB 模式（去除 alpha 通道）
    if image.mode in ('RGBA', 'LA', 'P'):
        rgb_image = Image.new('RGB', image.size, (255, 255, 255))
        if image.mode == 'P':
            image = image.convert('RGBA')
        rgb_image.paste(image, mask=image.split()[-1] if image.mode in ('RGBA', 'LA') else None)
        image = rgb_image
    
    image.save(buffered, format="JPEG", quality=85)
    img_bytes = buffered.getvalue()
    return base64.b64encode(img_bytes).decode('utf-8')


def _call_vision_api_with_fallback(base64_image: str, prompt: str) -> str:
    """
    使用多模型回退策略调用 Vision API
    委托给 llm_client.call_vision
    """
    try:
        return call_vision(base64_image, prompt)
    except Exception as e:
        return f"[图片页，提取失败：{str(e)}]"


def _extract_slide_text(slide) -> str:
    """从幻灯片中提取所有文本"""
    text_parts = []
    
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            text_parts.append(shape.text.strip())
    
    return "\n".join(text_parts)


def _slide_to_image(slide, slide_width, slide_height) -> Optional[Image.Image]:
    """
    将幻灯片转换为图片
    注意：python-pptx 不直接支持渲染幻灯片为图片
    这里我们尝试提取幻灯片中的图片，如果有的话
    """
    # 尝试找到幻灯片中的图片
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image_stream = io.BytesIO(shape.image.blob)
                return Image.open(image_stream)
            except Exception:
                continue
    
    return None


def _is_image_slide(slide) -> bool:
    """判断幻灯片是否为图片型（文字少于50个字符）"""
    text = _extract_slide_text(slide)
    return len(text.strip()) < 50


def _convert_pptx(file_path: str) -> str:
    """
    转换 PPTX 文件为 Markdown
    支持文字型和图片型幻灯片
    """
    try:
        prs = Presentation(file_path)
        markdown_parts = [f"# {Path(file_path).stem}\n"]
        
        for idx, slide in enumerate(prs.slides, start=1):
            # 提取幻灯片标题
            title = ""
            if slide.shapes.title:
                title = slide.shapes.title.text.strip()
            
            slide_header = f"## 第{idx}页：{title if title else '无标题'}\n"
            markdown_parts.append(slide_header)
            
            # 判断是文字型还是图片型
            if _is_image_slide(slide):
                # 图片型幻灯片：使用 Vision API
                print(f"  处理图片型幻灯片 {idx}/{len(prs.slides)}...")
                
                # 尝试提取图片
                slide_image = _slide_to_image(slide, prs.slide_width, prs.slide_height)
                
                if slide_image:
                    # 转换为 base64
                    base64_image = _image_to_base64(slide_image)
                    
                    # 调用 Vision API
                    prompt = """这是一张PPT幻灯片。请提取：
1) 标题（如果有）
2) 图表显示的数据、趋势和关键结论
3) 所有可见的文字标签和数值

用简洁的中文结构化文本输出，不要描述颜色、字体等视觉风格。"""
                    
                    vision_result = _call_vision_api_with_fallback(base64_image, prompt)
                    markdown_parts.append(vision_result + "\n")
                else:
                    # 没有找到图片，但文字很少
                    text = _extract_slide_text(slide)
                    if text:
                        markdown_parts.append(text + "\n")
                    else:
                        markdown_parts.append("[图片页，无法提取内容]\n")
            else:
                # 文字型幻灯片：直接提取文本
                text = _extract_slide_text(slide)
                markdown_parts.append(text + "\n")
        
        return "\n".join(markdown_parts)
        
    except Exception as e:
        raise PPTXConversionError(f"PPTX 转换失败: {str(e)}")


def _convert_pdf(file_path: str) -> str:
    """
    转换 PDF 文件为 Markdown
    自动判断是原生 PDF 还是扫描件
    """
    try:
        # Step 1: 尝试用 PyMuPDF 提取文本
        doc = fitz.open(file_path)
        
        # 提取前3页文本，统计字符数
        sample_text = ""
        pages_to_sample = min(3, len(doc))
        
        for page_num in range(pages_to_sample):
            page = doc[page_num]
            sample_text += page.get_text()
        
        doc.close()
        
        char_count = len(sample_text.strip())
        
        # Step 2: 判断是否为原生 PDF
        if char_count > 500:
            # 原生 PDF，使用 pymupdf4llm
            print(f"  检测到原生 PDF，使用 pymupdf4llm 提取...")
            try:
                import pymupdf4llm
                markdown_text = pymupdf4llm.to_markdown(file_path)
                return markdown_text
            except ImportError:
                raise PDFConversionError("pymupdf4llm 未安装")
        
        # Step 3: 扫描件或图片型 PDF，使用多种 OCR 方案
        print(f"  检测到扫描件/图片型 PDF（字符数: {char_count}），尝试 OCR 处理...")
        
        # 方案1：尝试 marker-pdf（最佳精度，但下载可能有问题）
        try:
            from marker.converters.pdf import PdfConverter
            from marker.models import create_model_dict
            from marker.output import text_from_rendered
            
            print(f"  使用 marker-pdf 转换器...")
            converter = PdfConverter(artifact_dict=create_model_dict())
            rendered = converter(file_path)
            text, _, _ = text_from_rendered(rendered)
            return text
            
        except Exception as e:
            error_msg = str(e)
            print(f"  ⚠ marker-pdf 失败: {error_msg[:100]}")
            
            # 如果是模型下载问题，尝试方案2
            if any(keyword in error_msg.lower() for keyword in ['download', 'model', 'offline', 'connect']):
                print(f"  → 切换到备用 OCR 方案...")
            else:
                # 其他错误直接抛出
                raise PDFConversionError(f"marker-pdf 处理失败: {error_msg}")
        
        # 方案2：pytesseract + pdf2image（备用方案）
        try:
            print(f"  尝试使用 pytesseract OCR...")
            from pdf2image import convert_from_path
            import pytesseract
            
            # 转换 PDF 为图片
            images = convert_from_path(file_path, dpi=200)
            
            # OCR 每一页
            text_parts = [f"# {Path(file_path).stem}\n"]
            for idx, image in enumerate(images, start=1):
                print(f"    处理第 {idx}/{len(images)} 页...")
                page_text = pytesseract.image_to_string(image, lang='chi_sim+eng')
                text_parts.append(f"\n## 第 {idx} 页\n")
                text_parts.append(page_text)
            
            return "\n".join(text_parts)
            
        except ImportError as ie:
            missing_lib = "pdf2image" if "pdf2image" in str(ie) else "pytesseract"
            raise PDFConversionError(
                f"OCR 库未安装。请选择以下方案之一:\n"
                f"1. 安装 marker-pdf: pip install marker-pdf\n"
                f"2. 安装 pytesseract: pip install pytesseract pdf2image\n"
                f"   (需要先安装 Tesseract-OCR: https://github.com/UB-Mannheim/tesseract/wiki)\n"
                f"缺失库: {missing_lib}"
            )
        except Exception as e:
            raise PDFConversionError(f"PDF OCR 失败: {str(e)}")
    
    except PDFConversionError:
        raise
    except Exception as e:
        raise PDFConversionError(f"PDF 转换失败: {str(e)}")


def _convert_docx(file_path: str) -> str:
    """使用 markitdown 转换 DOCX"""
    if MarkItDown is None:
        raise ConversionError("markitdown 未安装")
    
    try:
        md = MarkItDown()
        result = md.convert(file_path)
        return result.text_content
    except Exception as e:
        raise ConversionError(f"DOCX 转换失败: {str(e)}")


def _convert_html(file_path: str) -> str:
    """使用 markitdown 转换 HTML"""
    if MarkItDown is None:
        raise ConversionError("markitdown 未安装")
    
    try:
        md = MarkItDown()
        result = md.convert(file_path)
        return result.text_content
    except Exception as e:
        raise ConversionError(f"HTML 转换失败: {str(e)}")


def _convert_text(file_path: str) -> str:
    """直接读取文本文件"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        # 尝试其他编码
        try:
            with open(file_path, 'r', encoding='gbk') as f:
                return f.read()
        except Exception as e:
            raise ConversionError(f"文本文件读取失败: {str(e)}")
    except Exception as e:
        raise ConversionError(f"文本文件读取失败: {str(e)}")


def _convert_markdown_with_images(file_path: str) -> str:
    """
    转换 Markdown 文件，并处理其中的图片引用
    使用 Vision API 提取图片内容
    """
    import re
    
    # 读取 Markdown 内容
    content = _convert_text(file_path)
    
    # 查找所有图片引用：![alt](path) 或 ![alt](url)
    image_pattern = r'!\[([^\]]*)\]\(([^\)]+)\)'
    matches = list(re.finditer(image_pattern, content))
    
    if not matches:
        # 没有图片，直接返回
        return content
    
    print(f"  发现 {len(matches)} 张图片，开始处理...")
    
    # 处理每张图片
    base_dir = Path(file_path).parent
    result_parts = []
    last_end = 0
    
    for idx, match in enumerate(matches, start=1):
        # 添加图片前的文本
        result_parts.append(content[last_end:match.start()])
        
        alt_text = match.group(1)
        image_path = match.group(2)
        
        # 跳过网络图片（暂不支持）
        if image_path.startswith(('http://', 'https://')):
            result_parts.append(match.group(0))  # 保持原样
            result_parts.append(f"\n[网络图片，未处理]\n")
            last_end = match.end()
            continue
        
        # 解析本地图片路径
        full_image_path = base_dir / image_path
        
        if not full_image_path.exists():
            result_parts.append(match.group(0))
            result_parts.append(f"\n[图片不存在: {image_path}]\n")
            last_end = match.end()
            continue
        
        # 读取并转换图片
        try:
            print(f"    处理图片 {idx}/{len(matches)}: {image_path}")
            image = Image.open(full_image_path)
            base64_image = _image_to_base64(image)
            
            # 调用 Vision API
            prompt = f"""这是一张图片（描述：{alt_text if alt_text else '无'}）。
请提取图片中的所有信息：
1) 如果是图表，提取数据、趋势和结论
2) 如果是截图，提取所有可见文字
3) 如果是其他内容，简要描述

用简洁的中文输出，保持结构化。"""
            
            vision_result = _call_vision_api_with_fallback(base64_image, prompt)
            
            # 替换为提取的内容
            result_parts.append(f"\n**[图片内容: {alt_text if alt_text else image_path}]**\n")
            result_parts.append(vision_result)
            result_parts.append("\n")
            
        except Exception as e:
            result_parts.append(match.group(0))
            result_parts.append(f"\n[图片处理失败: {str(e)}]\n")
        
        last_end = match.end()
    
    # 添加剩余文本
    result_parts.append(content[last_end:])
    
    return "".join(result_parts)


def _convert_latex(file_path: str) -> str:
    """使用 pandoc 转换 LaTeX"""
    try:
        result = subprocess.run(
            ['pandoc', file_path, '-f', 'latex', '-t', 'markdown'],
            capture_output=True,
            text=True,
            timeout=60
        )
        
        if result.returncode != 0:
            raise ConversionError(f"pandoc 转换失败: {result.stderr}")
        
        return result.stdout
        
    except FileNotFoundError:
        raise ConversionError("pandoc 未安装，请先安装 pandoc")
    except subprocess.TimeoutExpired:
        raise ConversionError("LaTeX 转换超时")
    except Exception as e:
        raise ConversionError(f"LaTeX 转换失败: {str(e)}")


def convert_to_markdown(file_path: str) -> str:
    """
    主入口：将文件转换为 Markdown
    
    Args:
        file_path: 文件路径
        
    Returns:
        Markdown 格式的文本内容
        
    Raises:
        ConversionError: 转换失败时抛出
    """
    path = Path(file_path)
    
    if not path.exists():
        raise ConversionError(f"文件不存在: {file_path}")
    
    ext = path.suffix.lower()
    
    print(f"转换文件: {path.name} ({ext})")
    
    # 路由到对应的转换函数
    if ext == '.pptx':
        return _convert_pptx(file_path)
    elif ext == '.pdf':
        return _convert_pdf(file_path)
    elif ext == '.docx':
        return _convert_docx(file_path)
    elif ext in ['.html', '.htm']:
        return _convert_html(file_path)
    elif ext == '.md':
        # Markdown 文件：处理图片引用
        return _convert_markdown_with_images(file_path)
    elif ext == '.txt':
        return _convert_text(file_path)
    elif ext == '.tex':
        return _convert_latex(file_path)
    else:
        raise ConversionError(f"不支持的文件类型: {ext}")


if __name__ == "__main__":
    # 测试代码
    import sys
    
    if len(sys.argv) < 2:
        print("用法: python convert.py <文件路径>")
        sys.exit(1)
    
    try:
        markdown = convert_to_markdown(sys.argv[1])
        print("\n" + "="*50)
        print("转换结果:")
        print("="*50)
        print(markdown[:1000])  # 只打印前1000个字符
        if len(markdown) > 1000:
            print(f"\n... (还有 {len(markdown) - 1000} 个字符)")
    except ConversionError as e:
        print(f"错误: {e}")
        sys.exit(1)
